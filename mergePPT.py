#!/usr/bin/env python3
import sys
import os
import subprocess
import io

# --- Auto-install required modules ---
try:
    from pptx import Presentation
    from copy import deepcopy
except ImportError:
    print("Module 'python-pptx' not found. Attempting to install it now...")
    try:
        # Use pip to install the missing module
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        print("Successfully installed 'python-pptx'.")
        # Retry importing after installation
        from pptx import Presentation
        from copy import deepcopy
    except Exception as e:
        print(f"🚨 Failed to install 'python-pptx'. Please install it manually using 'pip install python-pptx'.")
        print(f"Error: {e}")
        sys.exit(1) # Exit the script if installation fails

def merge_presentations(ppt_files, output_filename):
    """
    Merges a list of PowerPoint files into a single file.
    The first presentation's slide master is used for all slides.
    """
    # Check if there are any files to merge
    if not ppt_files:
        print("⚠️ No PowerPoint files found to merge.")
        return

    print(f"Found {len(ppt_files)} presentations to merge.")

    # Use the first presentation in the list as the base
    base_ppt = Presentation(ppt_files[0])
    print(f"-> Using '{ppt_files[0]}' as the base presentation.")

    # Define the necessary XML namespaces
    ns = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }

    # Loop through the remaining presentations
    for ppt_file in ppt_files[1:]:
        try:
            print(f"-> Merging slides from '{ppt_file}'...")
            other_ppt = Presentation(ppt_file)

            # Iterate through each slide of the presentation to be merged
            for slide in other_ppt.slides:
                # Find the layout in the destination presentation that matches the source slide's layout.
                try:
                    layout_name = slide.slide_layout.name
                    dest_layout = next(l for l in base_ppt.slide_layouts if l.name == layout_name)
                except StopIteration:
                    # Fallback to a common layout (e.g., blank) if no match is found.
                    dest_layout = base_ppt.slide_layouts[6]

                new_slide = base_ppt.slides.add_slide(dest_layout)

                # Copy shapes from the old slide to the new one, using deepcopy for better fidelity.
                for shape in slide.shapes:
                    new_el = deepcopy(shape.element)
                    new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

                # Copy relationships for images and other media.
                for rel in slide.part.rels.values():
                    # Skip external relationships (like hyperlinks)
                    if rel.is_external:
                        continue
                    
                    # If it's an image, add it to the new presentation
                    if "image" in rel.target_part.content_type:
                        target_part = rel.target_part
                        image_stream = io.BytesIO(target_part.blob)
                        image_part, rId = new_slide.part.get_or_add_image_part(image_stream)
                        
                        # Find and update the relationship ID in the copied shapes
                        for shape_elm in new_slide.shapes.element.findall('.//a:blip', namespaces=ns):
                            # **FIXED LINE**: Use .get() for namespaced attributes
                            embed_rId = shape_elm.get(f"{{{ns['r']}}}embed")
                            if embed_rId == rel.rId:
                                # **FIXED LINE**: Use .set() to update the namespaced attribute
                                shape_elm.set(f"{{{ns['r']}}}embed", rId)
                                break

        except Exception as e:
            print(f"  -> 🚨 Error processing '{ppt_file}': {e}")

    # Save the final merged presentation
    base_ppt.save(output_filename)
    print(f"\n✅ Success! All presentations have been merged into '{output_filename}'.")


def main():
    """
    Main function to determine which files to merge.
    """
    output_filename = "mergedPPT.pptx"
    ppt_files_to_merge = []

    # Scenario 1: Specific files are provided as command-line arguments
    if len(sys.argv) > 1:
        print("Merging specific files listed as arguments...")
        # Filter for existing files ending with .pptx
        ppt_files_to_merge = [f for f in sys.argv[1:] if f.endswith('.pptx') and os.path.exists(f)]

    # Scenario 2: No arguments are given, so merge all pptx files in the directory
    else:
        print("No files specified. Merging all .pptx files in the current directory...")
        current_directory = os.getcwd()
        files = os.listdir(current_directory)
        # Filter for .pptx files, sort them alphabetically, and exclude the output file itself
        ppt_files_to_merge = sorted([f for f in files if f.endswith('.pptx') and f != output_filename])

    # Run the merge function with the determined list of files
    merge_presentations(ppt_files_to_merge, output_filename)


if __name__ == "__main__":
    main()

